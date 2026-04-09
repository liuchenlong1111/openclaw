#!/usr/bin/env python3
"""
Extract text content from PowerPoint presentations (.pptx).
"""

import sys
import argparse

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

def extract_pptx(file_path):
    """Extract text from PowerPoint presentation"""
    try:
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"=== Slide {slide_num} ==="]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting PPTX: {e}", file=sys.stderr)
        sys.exit(1)

def main():
    parser = argparse.ArgumentParser(description='Extract text from PowerPoint presentations')
    parser.add_argument('file_path', help='Path to .pptx file')
    parser.add_argument('--output', '-o', help='Output file (default: stdout)')
    args = parser.parse_args()
    
    text = extract_pptx(args.file_path)
    
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"Extracted text saved to {args.output}", file=sys.stderr)
    else:
        print(text)

if __name__ == '__main__':
    main()
